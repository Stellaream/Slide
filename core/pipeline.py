# core/pipeline.py
import os
import concurrent.futures
from pptx import Presentation
from engine.renderer import ProRenderer
from engine.size import fix_ppt_with_drag_simulation
from engine.image_manager import GlobalImageMatcher, StockManager

from config import OUTPUT_DIR, STOCK_DIR 
from utils import save_debug_file, get_random_background, extract_elements_robust, calculate_overlap, calculate_alignment, calculate_layout_score
from core.content import docx_to_markdown, collect_ref_chunks, collect_ref_images
from core.llm import get_ppt_outline, generate_single_slide

def run_pipeline(docx_path, log_callback=None, user_assets=None):
    """
    全自动化 PPT 生成流水线
    :param docx_path: 输入文档路径
    :param log_callback: (可选) 回调函数 func(msg, type)，用于向前端推送实时日志
    :return: 最终生成的 PPT 文件路径
    """

    # --- 内部辅助函数：双向日志 ---
    # 既打印到后台控制台，也推送到前端
    def log(msg, log_type="info"):
        print(f"[{log_type.upper()}] {msg}")
        if log_callback:
            log_callback(msg, log_type)

    log("启动全自动化 PPT 生成流程...", "info")

    # 0. 供 LLM 预感知的用户图片信息（语义+尺寸）
    user_asset_hints = []
    if user_assets:
        for idx, asset in enumerate(user_assets):
            user_asset_hints.append({
                "asset_id": f"I{idx + 1}",
                "tags": asset.get("tags", []),
                "aspect_ratio": asset.get("aspect_ratio")
            })
        log(f"检测到 {len(user_asset_hints)} 张用户图片，将参与大纲规划...", "info")
    
    # 1. 文档解析
    log("正在解析原始文档...", "info")
    chunks = docx_to_markdown(docx_path)
    if not chunks:
        log("❌ 文档解析失败，内容为空", "error")
        return None
    save_debug_file("source_chunks.json", chunks)


    # 2. 生成大纲
    log("AI 正在深度阅读文档并规划大纲 ...", "info")
    outline = get_ppt_outline(chunks, user_asset_hints)
    if not outline:
        log("❌ 大纲生成为空，流程终止。", "error")
        return None
    
    log(f"✅ 大纲规划完成，共 {len(outline)} 页。准备生成详细内容...", "success")
    save_debug_file("outline.json", outline)
    
    # 3. 准备背景
    bg_path = get_random_background()


    # 4. 并行生成详细页面 (每页生成2次选优)
    results = [None] * len(outline)       # 存放每次的最优结果
    worst_results = [None] * len(outline) # 存放每次的较差结果，用于对比
    
    completed_slides = 0 
    
    # 记录每页生成成功的候选结果: { idx:[res1, res2] }
    slide_candidates = {item["index"]:[] for item in outline}
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        future_to_info = {}
        
        for item in outline:
            idx = item["index"]

            # 抽取该页引用的原文与图片
            content = collect_ref_chunks(item, chunks)
            page_images = collect_ref_images(item, user_asset_hints)

            # 为每一页提交 2 次生成任务
            for candidate_id in range(2):
                future = executor.submit(
                    generate_single_slide,
                    item,
                    content,
                    page_images
                )
                # 记录这个 future 属于哪一页的第几次生成
                future_to_info[future] = {"index": idx, "candidate_id": candidate_id}
        
        # 收集结果
        for future in concurrent.futures.as_completed(future_to_info):
            info = future_to_info[future]
            idx = info["index"]
            
            try:
                res = future.result()
                slide_candidates[idx].append(res)
            except Exception as e:
                log(f"⚠️ 第 {idx} 页 (候选 {info['candidate_id']}) 生成线程异常: {e}", "warning")
                # 即使失败也塞入一个 None 占位，以便后续判断该页两次任务都已结束
                slide_candidates[idx].append(None)
            
            # 当该页的 2 次任务全部执行完毕时，进行打分比较
            if len(slide_candidates[idx]) == 2:
                # 过滤掉生成失败的 None 结果
                valid_candidates =[c for c in slide_candidates[idx] if c is not None]
                
                best_res = None
                worst_res = None
                
                if len(valid_candidates) == 2:
                    # 两次都成功，计算得分并比较
                    score_list =[]
                    for c in valid_candidates:
                        elements = c.get("elements",[])
                        overlap_err = calculate_overlap(elements)
                        align_err = calculate_alignment(elements)
                        score = calculate_layout_score(overlap_err, align_err)
                        score_list.append(score)
                    
                    # 取分数高的作为 best，低的作为 worst
                    if score_list[0] >= score_list[1]:
                        best_res, worst_res = valid_candidates[0], valid_candidates[1]
                        best_res["_layout_score"] = score_list[0] # 可选：将分数记录在 JSON 中方便你调试看
                        worst_res["_layout_score"] = score_list[1]
                    else:
                        best_res, worst_res = valid_candidates[1], valid_candidates[0]
                        best_res["_layout_score"] = score_list[1]
                        worst_res["_layout_score"] = score_list[0]
                        
                elif len(valid_candidates) == 1:
                    best_res = valid_candidates[0]
                    worst_res = None 
                else:
                    log(f"❌ 第 {idx} 页两次生成均失败！", "error")
                
                # 按照原始索引存入对应的数组位
                results[idx - 1] = best_res
                worst_results[idx - 1] = worst_res
                
                completed_slides += 1
                log(f"正在构建页面内容与布局并择优... ({completed_slides}/{len(outline)})", "info")

    # 分别保存最优结果与淘汰结果
    save_debug_file("merged.json", results) 
    save_debug_file("merged_worst.json", worst_results)

    # 5. 图片资源的智能处理
    log("正在进行图片资源的全局匹配...", "info")
    
    matcher = GlobalImageMatcher(user_assets)
    # 得到匹配字典: { (页码idx, 元素idx): "path/to/img.jpg" }
    mapping_result = matcher.run_matching(results)

    stock_mgr = None
    if os.path.exists(STOCK_DIR):
        print(f"加载本地图库: {STOCK_DIR}") 
        stock_mgr = StockManager(assets_dir=STOCK_DIR)
    else:
        print(f"⚠️ 未找到 stock 文件夹，兜底功能失效")
        stock_mgr = StockManager(assets_dir=None)
    
    # 6. 渲染阶段
    log("正在启动渲染引擎...", "info")
    prs = Presentation()
    renderer = ProRenderer(prs)
    
    for page_idx, slide_data in enumerate(results):
        elements = extract_elements_robust(slide_data)
        if not elements: continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        for el_idx, el in enumerate(elements):
            final_img_path = None
            
            # 优先用 AI 分配的用户图
            if (page_idx, el_idx) in mapping_result:
                final_img_path = mapping_result[(page_idx, el_idx)]
            
            # 没图则用本地库兜底
            elif el.get('type') == 'image' and stock_mgr:
                final_img_path = stock_mgr.pick_next()

            try:
                renderer.render_element(slide, el, force_image_path=final_img_path)
            except Exception as e:
                print(f"渲染异常: {e}")

    # 7. 背景与保存
    if bg_path:
        log("正在应用全局视觉风格...", "info")
        renderer.add_background_to_all_slides(bg_path)

    origin_path = os.path.join(OUTPUT_DIR, "Origin.pptx")
    prs.save(origin_path)
    
    # 8. 字号修复
    if os.path.exists(origin_path):
        log("正在启动 PowerPoint 引擎进行版式自适应修复...", "info")
        try:
            # 调用 engine/size.py 中的逻辑
            fix_ppt_with_drag_simulation(origin_path)
            log("✅ 全流程完成！最终文件已生成。", "success")
        except Exception as e:
            log(f"⚠️ 字号修复失败 (可能需要 Windows 环境): {e}", "warning")
            
    return origin_path