import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

# ==========================================
# 1. 配色主题配置 (Academic Blue)
# ==========================================
THEME = {
    "primary":    "#005691",  # 主色：深蓝
    "accent":     "#1E88E5",  # 辅色：亮蓝
    "bg_card":    "#FFFFFF",  # 卡片背景
    "border":     "#D1E1EF",  # 边框色
    "text_main":  "#333333",  # 正文
    "text_sub":   "#5F6368",  # 辅助文字
    "line":       "#E0E0E0"   # 分隔线
}

class ProRenderer:
    def __init__(self, prs):
        self.prs = prs
        # 16:9 宽屏设置
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # 栅格系统
        self.cols = 16 
        self.rows = 9
        self.grid_w = self.prs.slide_width / self.cols
        self.grid_h = self.prs.slide_height / self.rows

    def _hex_to_rgb(self, hex_str):
        if not hex_str or str(hex_str).lower() in ['transparent', 'none', '']: 
            return None
        try:
            hex_str = str(hex_str).lstrip('#')
            if len(hex_str) == 3: 
                hex_str = ''.join([c*2 for c in hex_str])
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        except: 
            return RGBColor(0, 0, 0)

    def _set_font_face(self, run, font_name='Microsoft YaHei'):
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        ea = OxmlElement('a:ea')
        ea.set('typeface', font_name)
        rPr.append(ea)

    def _apply_text_frame_style(self, shape, content, align='left', font_size=18, color='#333333', bold=False, line_spacing=1.2, margin=5, auto_size=True):
        """ 
        style应用函数 
        :param auto_size: True=开启自动缩放(font_size作为上限), False=强制固定字号(font_size作为绝对值)
        """
        tf = shape.text_frame
        tf.clear()
        
        # 设置内边距
        tf.margin_left = Pt(margin)
        tf.margin_right = Pt(margin)
        tf.margin_top = Pt(margin)
        tf.margin_bottom = Pt(margin)
        
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE 
        tf.word_wrap = True 
        
        # 根据 auto_size 参数决定策略
        if auto_size:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        else:
            tf.auto_size = MSO_AUTO_SIZE.NONE # 关闭自动调整，严格执行字号

        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        content_str = str(content).strip()
        lines = content_str.split('\n') if content_str else [""]
        rgb_color = self._hex_to_rgb(color)

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
            p.line_spacing = line_spacing
            p.space_after = Pt(6) if (len(lines) > 1 and i < len(lines)-1) else Pt(0)

            run = p.add_run()
            run.text = line.strip()
            run.font.size = Pt(font_size) # 若 auto_size=True，这只是最大值；若 False，这是固定值
            run.font.bold = bold
            if rgb_color:
                run.font.color.rgb = rgb_color
            
            self._set_font_face(run, 'Microsoft YaHei')

    def render_element(self, slide, el, force_image_path=None):
        pos = el.get('pos', {})
        gx = max(0, min(float(pos.get('x', 0)), 16))
        gy = max(0, min(float(pos.get('y', 0)), 9))
        gw = max(0.5, min(float(pos.get('w', 4)), 16 - gx))
        gh = max(0.5, min(float(pos.get('h', 2)), 9 - gy))
        
        l = int(gx * self.grid_w)
        t = int(gy * self.grid_h)
        w = int(gw * self.grid_w)
        h = int(gh * self.grid_h)
        
        el_type = str(el.get('type', 'text')).lower().strip()
        content = el.get('content', '')
        subtitle = el.get('subtitle', '')
        style = el.get('style', {})

        # 1. Title (大标题 - 保持醒目风格)
        if el_type == 'title':
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
            shape.adjustments[0] = 0.2
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
            
            shape.line.color.rgb = self._hex_to_rgb(THEME['primary'])
            shape.line.width = Pt(2.0)

            # 标题通常文字较少，建议默认开启 AutoSize 防止溢出，但也允许手动覆盖（这里暂定默认Auto）
            self._apply_text_frame_style(
                shape, content, 
                align='center', 
                font_size=40, 
                color=THEME['primary'], 
                bold=True,
                auto_size=True 
            )

        # 2. Text (普通文本 - 智能判定 Explicit vs Auto)
        elif el_type == 'text':
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
            
            # 样式处理
            bg_color = style.get('bg_color', 'transparent')
            if bg_color != 'transparent':
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._hex_to_rgb(bg_color)
            else:
                shape.fill.background()

            if style.get('border'):
                shape.line.color.rgb = self._hex_to_rgb(THEME['border'])
                shape.line.width = Pt(1)
            else:
                shape.line.fill.background()

            json_align = style.get('align', 'left')
            is_bold = style.get('bold', False)
            text_color = style.get('color', THEME['primary'] if is_bold else THEME['text_main'])
            
            # 检查 font_size 是否设置
            specified_fs = style.get('font_size')
            
            if specified_fs:
                # 显式设置了字号 -> 关闭自动缩放，使用指定值
                use_auto_size = False
                fs_val = specified_fs
            else:
                # 未设置字号 -> 开启自动缩放，使用默认最大值
                use_auto_size = True
                fs_val = 24 # 默认最大字号

            self._apply_text_frame_style(
                shape, content, 
                align=json_align, 
                font_size=fs_val,  
                color=text_color, 
                bold=is_bold,
                auto_size=use_auto_size
            )

        # 3. Card
        elif el_type == 'card':
            # 背景容器
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
            bg_shape.adjustments[0] = 0.04
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self._hex_to_rgb(THEME['bg_card'])
            bg_shape.line.color.rgb = self._hex_to_rgb(THEME['border'])
            bg_shape.line.width = Pt(1.5)

            # 布局计算
            header_h = max(int(Inches(0.5)), min(int(h * 0.25), int(Inches(0.9))))
            body_h = h - header_h

            # Header
            header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, header_h)
            header_box.fill.background() 
            header_box.line.fill.background() 
            
            self._apply_text_frame_style(
                header_box, subtitle, 
                align='center', font_size=22, color=THEME['primary'], bold=True, margin=3,
                auto_size=True # 卡片标题建议自动适配
            )

            # Line
            line_y = t + header_h
            line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, l + int(w * 0.1), line_y, l + int(w * 0.9), line_y)
            line_shape.line.color.rgb = self._hex_to_rgb(THEME['line'])
            line_shape.line.width = Pt(1)

            # Body
            body_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, line_y, w, body_h)
            body_box.fill.background()
            body_box.line.fill.background()
            
            # 卡片内容通常由 AutoSize 兜底，防止溢出
            self._apply_text_frame_style(
                body_box, content, 
                align='center', font_size=16, color=THEME['text_main'], bold=False, margin=8,
                auto_size=True
            )

        # 4. Image
        elif el_type == 'image':
            img_path = force_image_path
            
            if img_path and os.path.exists(img_path):
                try:
                    pic = slide.shapes.add_picture(img_path, l, t)
                    self._crop_image_to_fit(pic, w, h)
                    pic.line.color.rgb = self._hex_to_rgb(THEME['border'])
                    pic.line.width = Pt(1)
                    return 
                except: pass

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
            shape.adjustments[0] = 0.02
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(248, 248, 248)
            shape.line.color.rgb = self._hex_to_rgb(THEME['border'])
            shape.line.dash_style = 4
            shape.line.width = Pt(1.5)
            
            self._apply_text_frame_style(
                shape, f"🖼️\n{content}", 
                align="center", font_size=14, color=THEME['text_sub'], bold=False,
                auto_size=True
            )

    def _crop_image_to_fit(self, pic, target_w, target_h):
        img_ratio = pic.width / pic.height
        target_ratio = target_w / target_h
        
        pic.crop_left = 0
        pic.crop_right = 0
        pic.crop_top = 0
        pic.crop_bottom = 0

        if img_ratio > target_ratio:
            new_width = target_h * img_ratio
            crop_amount = (new_width - target_w) / new_width
            pic.crop_left = crop_amount / 2
            pic.crop_right = crop_amount / 2
            pic.width = target_w
            pic.height = target_h
        else:
            new_height = target_w / img_ratio
            crop_amount = (new_height - target_h) / new_height
            pic.crop_top = crop_amount / 2
            pic.crop_bottom = crop_amount / 2
            pic.width = target_w
            pic.height = target_h

    def add_background_to_all_slides(self, image_path):
        if not image_path or not os.path.exists(image_path):
            return

        for slide in self.prs.slides:
            try:
                pic = slide.shapes.add_picture(
                    image_path, 0, 0,
                    width=self.prs.slide_width,
                    height=self.prs.slide_height
                )

                # 放到底层
                spTree = slide.shapes._spTree
                element = pic._element
                spTree.remove(element)
                spTree.insert(2, element)

            except Exception as e:
                pass
